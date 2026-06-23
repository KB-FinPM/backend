# EN: Exports generated artifact JSON to user-facing files and uploads them.
# KO: 생성 산출물 JSON을 사용자 파일로 변환하고 S3에 업로드합니다.

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from datetime import date
from io import BytesIO
from pathlib import PurePath
from typing import Any
from uuid import uuid4

from app.core.config import settings
from app.core.logger import get_logger
from app.schemas.artifact import ArtifactType, DocumentMetadata, DocumentType
from app.services.document_service import DocumentService
from app.storage.s3 import S3Service, s3_service
from util.agent_template_utils import (
    build_placeholder_values,
    build_template_context,
    get_nested,
    get_value,
    load_output_mapper,
    output_file_name,
    resolve_template_path,
)

logger = get_logger(__name__)


@dataclass(frozen=True)
class ArtifactExportResult:
    storage_path: str
    file_name: str
    content_type: str
    document: DocumentMetadata | None = None


class ArtifactExportService:
    """Creates xlsx/pptx files for generated artifacts.

    Requirement specs are also registered as generated REQUIREMENT_SPEC documents
    so WBS and screen-design generation can use the returned document_id as their
    source document.
    """

    async def export_artifact(
        self,
        *,
        project_id: str,
        artifact_id: str,
        artifact_type: ArtifactType,
        result_json: dict[str, Any],
        project_name: str | None = None,
        document_service: DocumentService | None = None,
        storage_service: S3Service | None = None,
    ) -> ArtifactExportResult | None:
        if artifact_type == ArtifactType.REQUIREMENT_SPEC:
            file_name = self._requirement_file_name()
            file_bytes = self._build_requirement_xlsx(result_json)
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif artifact_type == ArtifactType.UNITTEST_SPEC:
            file_name = output_file_name("unit_test", "단위테스트케이스.xlsx")
            file_bytes = self._build_unit_test_xlsx(result_json)
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif artifact_type == ArtifactType.WBS:
            file_name = output_file_name("wbs", "WBS.xlsx")
            file_bytes = self._build_wbs_xlsx(result_json)
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif artifact_type == ArtifactType.SCREEN_DESIGN:
            file_name = output_file_name("screen_plan", "화면기획서.pptx")
            file_bytes = self._build_screen_design_pptx(result_json)
            content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            return None

        safe_file_name = PurePath(file_name).name
        storage_key = (
            f"{settings.S3_GENERATED_PREFIX}/{project_id}/"
            f"{artifact_type.value}/{artifact_id}/{safe_file_name}"
        )
        # Use the same storage service that handled the upload API whenever it is
        # available. This prevents generated artifacts from accidentally using a
        # different/global S3 backend configuration than uploaded source files.
        active_storage_service = (
            storage_service
            or (document_service.storage_service if document_service is not None else None)
            or s3_service
        )
        storage_path = await active_storage_service.upload(
            file_bytes=file_bytes,
            key=storage_key,
            content_type=content_type,
        )
        logger.info(
            "[ArtifactExport] exported | "
            f"project_id={project_id} | artifact_id={artifact_id} | "
            f"file_name={safe_file_name} | path={storage_path}"
        )

        generated_document = None
        if artifact_type == ArtifactType.REQUIREMENT_SPEC and document_service is not None:
            document_id = f"DOC-{uuid4().hex[:12].upper()}"
            generated_document = await document_service.ingest_uploaded_document(
                document_id=document_id,
                project_id=project_id,
                document_type=DocumentType.REQUIREMENT_SPEC,
                file_name=safe_file_name,
                storage_path=storage_path,
                file_bytes=file_bytes,
                parsed_context={
                    "text": self._requirement_text(result_json),
                    "requirements": result_json.get("requirements", []),
                    "metadata": {
                        "generated_from_artifact_id": artifact_id,
                        "artifact_type": artifact_type.value,
                        "source_file_name": safe_file_name,
                        "content_type": content_type,
                    },
                    "parser_name": "ArtifactExportService",
                },
            )
        elif artifact_type == ArtifactType.WBS and document_service is not None:
            document_id = f"DOC-{uuid4().hex[:12].upper()}"
            wbs_context = self._wbs_context(result_json, file_name=safe_file_name)
            generated_document = await document_service.ingest_uploaded_document(
                document_id=document_id,
                project_id=project_id,
                document_type=DocumentType.WBS,
                file_name=safe_file_name,
                storage_path=storage_path,
                file_bytes=file_bytes,
                parsed_context={
                    "text": self._wbs_text(result_json),
                    "metadata": {
                        "generated_from_artifact_id": artifact_id,
                        "artifact_type": artifact_type.value,
                        "source_file_name": safe_file_name,
                        "content_type": content_type,
                        "wbs_context": wbs_context,
                    },
                    "wbs_context": wbs_context,
                    "parser_name": "ArtifactExportService",
                },
            )

        return ArtifactExportResult(
            storage_path=storage_path,
            file_name=safe_file_name,
            content_type=content_type,
            document=generated_document,
        )

    def _requirement_file_name(self) -> str:
        return "요구사항명세서.xlsx"

    def _requirement_text(self, result_json: dict[str, Any]) -> str:
        lines = ["# REQUIREMENT_SPEC"]
        for item in result_json.get("requirements", []):
            metadata = item.get("metadata") or {}
            lines.append(
                " | ".join(
                    [
                        str(item.get("requirement_id", "")),
                        str(metadata.get("biz_requirement_name", "")),
                        str(metadata.get("domain", "")),
                        str(item.get("title", "")),
                        str(item.get("description", "")),
                    ]
                )
            )
        return "\n".join(lines)

    def _wbs_text(self, result_json: dict[str, Any]) -> str:
        lines = ["# WBS"]
        for task in result_json.get("tasks", []):
            if not isinstance(task, dict):
                continue
            source = self._wbs_source(task)
            lines.append(
                " | ".join(
                    str(value or "")
                    for value in [
                        source.get("task_id"),
                        source.get("id"),
                        source.get("wbs_name"),
                        source.get("start_date"),
                        source.get("end_date"),
                        source.get("worker"),
                        source.get("work_status") or source.get("status"),
                    ]
                )
            )
        return "\n".join(lines)

    def _wbs_context(
        self,
        result_json: dict[str, Any],
        *,
        file_name: str,
    ) -> dict[str, Any]:
        rows: list[dict[str, Any]] = []
        for row_number, task in enumerate(result_json.get("tasks", []), start=1):
            if not isinstance(task, dict):
                continue
            source = self._wbs_source(task)
            title = str(source.get("wbs_name") or source.get("name") or "").strip()
            if not title:
                continue
            start_date = source.get("start_date") or source.get("planned_start_date")
            end_date = source.get("end_date") or source.get("planned_end_date")
            worker = source.get("worker") or source.get("assignee") or source.get("owner")
            deliverable = source.get("deliverable") or source.get("artifact")
            status = (
                source.get("work_status")
                or source.get("status")
                or source.get("raw_status")
            )
            rows.append(
                {
                    **task,
                    "row_number": row_number,
                    "no": row_number,
                    "level": source.get("level"),
                    "레벨": source.get("level"),
                    "wbs_id": source.get("id"),
                    "ID": source.get("id"),
                    "title": title,
                    "WBS명": title,
                    "planned_start_date": start_date,
                    "시작예정일": start_date,
                    "planned_end_date": end_date,
                    "종료예정일": end_date,
                    "raw_assignee": worker,
                    "작업자": worker,
                    "artifact": deliverable,
                    "산출물": deliverable,
                    "raw_status": status,
                    "작업상태": status,
                    "source_document_name": file_name,
                }
            )

        return {
            "source_document_name": file_name,
            "sheet_name": "WBS",
            "header_row_number": 1,
            "columns": {
                "no": "NO",
                "level": "레벨",
                "id": "ID",
                "title": "WBS명",
                "planned_start_date": "시작예정일",
                "planned_end_date": "종료예정일",
                "assignee": "작업자",
                "artifact": "산출물",
                "status": "작업상태",
            },
            "rows": rows,
        }

    def _build_requirement_xlsx(self, result_json: dict[str, Any]) -> bytes:
        from openpyxl import Workbook, load_workbook
        from openpyxl.cell.cell import MergedCell

        mapper = load_output_mapper()
        req_mapper = get_nested(mapper, "requirement_spec", default={}) or {}
        sheet_cfg = req_mapper.get("data_sheet") or {}
        columns = self._legacy_requirement_columns(sheet_cfg.get("columns") or [])
        context = build_template_context(project_id=str((result_json.get("metadata") or {}).get("project_id") or ""), context=result_json.get("metadata") or {})

        template_path = req_mapper.get("template_path")
        if template_path:
            wb = load_workbook(resolve_template_path(template_path))
            self._apply_placeholder_sheets(wb, req_mapper.get("placeholder_sheets", []), context)
            target_sheet_name = sheet_cfg.get("sheet_name", "요구사항명세서")
            target_sheet_names = [target_sheet_name]
            # Existing requirement templates used by the team expose the same
            # requirement table through legacy sheets such as 감리제출용. Fill
            # them as well when present so the output matches pre-merge files.
            for legacy_name in ["감리제출용", "(요건분리 전)요구사항명세서"]:
                if legacy_name in wb.sheetnames and legacy_name not in target_sheet_names:
                    target_sheet_names.append(legacy_name)
            for sheet_name in target_sheet_names:
                if sheet_name not in wb.sheetnames:
                    continue
                ws = wb[sheet_name]
                self._write_mapped_rows(ws, result_json.get("requirements", []), sheet_cfg, columns, context, self._requirement_source)
                self._write_requirement_reference_fields(
                    ws,
                    result_json.get("requirements", []),
                    sheet_cfg,
                )
            return self._save_workbook(wb)

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_cfg.get("sheet_name") or "요구사항명세서"
        headers = [(column.get("header_names") or [column.get("field", "")])[0] for column in columns] or ["구분", "Biz요건ID", "Biz요건명", "요구사항ID", "요구사항명", "기능/비기능요구사항", "검토의견"]
        ws.append(headers)
        for row_number, item in enumerate(result_json.get("requirements", []), start=1):
            source = self._requirement_source(item)
            ws.append([get_value(source, column.get("field"), context=context, row_number=row_number) for column in columns] if columns else [
                source.get("category", ""), source.get("biz_requirement_id", ""), source.get("biz_requirement_name") or source.get("domain", ""), source.get("requirement_id", ""), source.get("requirement_name", ""), source.get("requirement_type", ""), source.get("note", ""),
            ])
        self._write_requirement_reference_fields(
            ws,
            result_json.get("requirements", []),
            sheet_cfg,
        )
        self._style_sheet(ws, widths=[18, 18, 24, 18, 36, 22, 60, 20, 20, 20, 20, 20, 20, 20, 20, 60])
        return self._save_workbook(wb)

    def _build_wbs_xlsx(self, result_json: dict[str, Any]) -> bytes:
        from openpyxl import Workbook, load_workbook

        mapper = load_output_mapper()
        wbs_mapper = get_nested(mapper, "wbs", default={}) or {}
        sheet_cfg = wbs_mapper.get("data_sheet") or {}
        columns = sheet_cfg.get("columns") or []
        context = build_template_context(project_id=str((result_json.get("metadata") or {}).get("project_id") or ""), context=result_json.get("metadata") or {})
        tasks = self._normalize_wbs_hierarchy(result_json.get("tasks", []), project_name=context.get("project_name") or "프로젝트명")
        self._apply_wbs_schedule_defaults(tasks, result_json)

        template_path = wbs_mapper.get("template_path")
        if template_path:
            wb = load_workbook(resolve_template_path(template_path))
            ws = wb[sheet_cfg.get("sheet_name", "WBS")]
            self._write_mapped_rows(ws, tasks, sheet_cfg, columns, context, self._wbs_source)
            self._write_wbs_display_ids(ws, tasks, sheet_cfg, columns)
            self._write_wbs_required_schedule_fields(ws, tasks, sheet_cfg)
            self._apply_wbs_visual_formatting(ws, tasks, sheet_cfg)
            return self._save_workbook(wb)

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_cfg.get("sheet_name") or "WBS"
        headers = [(column.get("header_names") or [column.get("field", "")])[0] for column in columns] or [
            "NO",
            "레벨",
            "ID",
            "WBS명",
            "시작예정일",
            "종료예정일",
            "작업자",
            "산출물",
        ]
        ws.append(headers)
        for row_number, task in enumerate(tasks, start=1):
            source = self._wbs_source(task)
            ws.append(
                [get_value(source, column.get("field"), context=context, row_number=row_number) for column in columns]
                if columns
                else [
                    row_number,
                    source.get("level", ""),
                    source.get("id", ""),
                    source.get("wbs_name", ""),
                    source.get("start_date", ""),
                    source.get("end_date", ""),
                    source.get("worker", ""),
                    source.get("deliverable", ""),
                ]
            )
        self._write_wbs_display_ids(ws, tasks, sheet_cfg, columns)
        self._write_wbs_required_schedule_fields(ws, tasks, sheet_cfg)
        self._style_sheet(ws, widths=[10, 10, 18, 48, 34, 20, 20, 34, 20])
        self._apply_wbs_visual_formatting(ws, tasks, sheet_cfg)
        return self._save_workbook(wb)

    def _build_unit_test_xlsx(self, result_json: dict[str, Any]) -> bytes:
        from openpyxl import Workbook, load_workbook

        mapper = load_output_mapper()
        unit_mapper = get_nested(mapper, "unit_test", default={}) or {}
        sheet_cfg = unit_mapper.get("data_sheet") or self._default_unit_test_sheet_cfg()
        columns = sheet_cfg.get("columns") or self._unit_test_columns()
        context = build_template_context(
            project_id=str((result_json.get("metadata") or {}).get("project_id") or ""),
            context=result_json.get("metadata") or {},
        )
        test_cases = result_json.get("test_cases", [])

        template_path = unit_mapper.get("template_path") or "template/탬플릿_단위테스트케이스.xlsx"
        if template_path:
            wb = load_workbook(resolve_template_path(template_path))
            self._apply_placeholder_sheets(
                wb,
                unit_mapper.get("placeholder_sheets") or self._default_cover_placeholder_sheets(),
                context,
            )
            sheet_name = self._resolve_sheet_name(
                wb,
                sheet_cfg.get("sheet_name", "케이스"),
                sheet_cfg.get("fallback_sheet_names") or ["목록총괄표"],
            )
            ws = wb[sheet_name]
            self._write_mapped_rows(ws, test_cases, sheet_cfg, columns, context, self._unit_test_source)
            self._format_unit_test_content_column(ws, sheet_cfg, columns, len(test_cases))
            return self._save_workbook(wb)

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_cfg.get("sheet_name") or "케이스"
        headers = [(column.get("header_names") or [column.get("field", "")])[0] for column in columns] or [
            "No",
            "요구사항ID",
            "요구사항명",
            "시나리오ID",
            "테스트케이스ID",
            "테스트케이스명",
            "테스트케이스 점검 내용",
        ]
        ws.append(headers)
        for row_number, item in enumerate(test_cases, start=1):
            source = self._unit_test_source(item)
            ws.append(
                [get_value(source, column.get("field"), context=context, row_number=row_number) for column in columns]
                if columns
                else [
                    row_number,
                    source.get("requirement_id", ""),
                    source.get("requirement_name", ""),
                    source.get("scenario_id", ""),
                    source.get("test_case_id", ""),
                    source.get("test_case_name", ""),
                    source.get("test_content", ""),
                ]
            )
        self._format_unit_test_content_column(ws, sheet_cfg, columns, len(test_cases))
        self._style_sheet(ws, widths=[8, 18, 32, 18, 22, 36, 72, 22, 18, 18])
        return self._save_workbook(wb)

    def _default_cover_placeholder_sheets(self) -> list[dict[str, Any]]:
        return [
            {
                "sheet_name": "표지",
                "placeholders": {
                    "{프로젝트명}": "project_name",
                    "{작성자명}": "author",
                    "{작성자}": "author",
                },
            },
            {
                "sheet_name": "개정이력",
                "placeholders": {
                    "{작성일}": "today",
                    "{작성자명}": "author",
                    "{작성자}": "author",
                },
            },
        ]

    def _default_unit_test_sheet_cfg(self) -> dict[str, Any]:
        return {
            "sheet_name": "케이스",
            "fallback_sheet_names": ["목록총괄표"],
            "header_row": 1,
            "start_row": 2,
            "clear_existing_rows": True,
            "copy_template_row_style": False,
            "columns": self._unit_test_columns(),
        }

    def _unit_test_columns(self) -> list[dict[str, Any]]:
        return [
            {"field": "row_number", "header_names": ["No", "NO"], "default_column": 1},
            {"field": "new_category", "header_names": ["신규구분"], "optional": True, "default_column": 2},
            {"field": "test_case_id", "header_names": ["테스트케이스ID", "테스트ID"], "default_column": 3},
            {"field": "test_case_name", "header_names": ["테스트케이스명", "화면명"], "default_column": 4},
            {"field": "scenario_id", "header_names": ["시나리오ID", "화면/배치 ID"], "default_column": 5},
            {"field": "screen_batch_type", "header_names": ["화면/배치"], "optional": True},
            {"field": "test_content", "header_names": ["테스트케이스 점검 내용", "  테스트케이스 점검 내용", "테스트케이스 내용"], "default_column": 7},
            {"field": "requirement_id", "header_names": ["요구사항 ID", "요구사항ID"], "optional": True},
            {"field": "requirement_name", "header_names": ["요구사항명"], "optional": True},
            {"field": "author", "header_names": ["담당 테스터", "담당자"], "optional": True, "default_column": 10},
        ]

    def _build_screen_design_pptx(self, result_json: dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        mapper = load_output_mapper()
        screen_mapper = get_nested(mapper, "screen_plan", default={}) or {}
        context = build_template_context(project_id=str((result_json.get("metadata") or {}).get("project_id") or ""), context=result_json.get("metadata") or {})
        template_path = screen_mapper.get("template_path")

        if template_path:
            prs = Presentation(resolve_template_path(template_path))
            common_values = self._common_screen_placeholder_values(screen_mapper, context)
            self._replace_placeholders_in_presentation(prs, common_values)
            template_slide_index = int(screen_mapper.get("template_slide_index", 2))
            screen_placeholders = get_nested(screen_mapper, "placeholder_slides", "screen_item", default={}) or {}
            table_mapper = screen_mapper.get("description_table") or {}
            for screen in result_json.get("screens", []):
                slide = self._duplicate_slide(prs, template_slide_index)
                source = self._screen_source(screen)
                values = {}
                values.update(common_values)
                values.update(build_placeholder_values(screen_placeholders, source=source, context=context))
                values.update(self._screen_placeholder_values(source, context))
                self._replace_placeholders_in_slide(slide, values)
                filled = self._fill_description_table(slide, source, table_mapper)
                if not filled:
                    self._add_description_fallback(slide, source)
                self._replace_placeholders_in_slide(slide, values)
            if template_slide_index < len(prs.slides):
                self._delete_slide(prs, template_slide_index)
            self._replace_placeholders_in_presentation(prs, common_values)
            bio = BytesIO()
            prs.save(bio)
            return bio.getvalue()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "화면설계서"
        title_slide.placeholders[1].text = "Generated by PM Artifact Agent"

        for screen in result_json.get("screens", []):
            source = self._screen_source(screen)
            slide = prs.slides.add_slide(self._get_blank_slide_layout(prs))
            title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.4), Inches(0.55))
            title_tf = title.text_frame
            title_tf.text = f"{source.get('screen_id', '')}  {source.get('screen_name', '')}"
            title_tf.paragraphs[0].font.size = Pt(20)
            title_tf.paragraphs[0].font.bold = True

            placeholder = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(7.5), Inches(5.95))
            placeholder.text = "화면 기획 내용 이미지 영역"
            placeholder.text_frame.paragraphs[0].font.size = Pt(18)

            rows = max(2, min(8, len(source.get("display_items", [])) + 1))
            table_shape = slide.shapes.add_table(rows, 2, Inches(8.25), Inches(1.05), Inches(4.65), Inches(5.95))
            table = table_shape.table
            table.cell(0, 0).text = "표시항목"
            table.cell(0, 1).text = "설명"
            for idx, item in enumerate(source.get("display_items", [])[: rows - 1], start=1):
                table.cell(idx, 0).text = str(item.get("item_name", ""))
                table.cell(idx, 1).text = str(item.get("description", ""))

            note = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35))
            note.text_frame.text = f"연관 요구사항: {source.get('requirement_id', '')}"
            note.text_frame.paragraphs[0].font.size = Pt(10)

        bio = BytesIO()
        prs.save(bio)
        return bio.getvalue()


    def _legacy_requirement_columns(self, columns: list[dict[str, Any]]) -> list[dict[str, Any]]:
        legacy_columns = [
            {"field": "biz_requirement_id", "header_names": ["Biz요건ID", "Biz 요건 ID"], "default_column": 1},
            {"field": "biz_requirement_name|domain", "header_names": ["Biz요건명", "업무", "업무명"], "default_column": 2},
            {"field": "category", "header_names": ["요구사항구분", "구분"], "default_column": 3},
            {"field": "creation_stage", "header_names": ["생성단계"], "optional": True, "default_column": 4},
            {"field": "status", "header_names": ["상태구분", "상태"], "optional": True, "default_column": 5},
            {"field": "source", "header_names": ["출처"], "optional": True, "default_column": 6},
            {"field": "requirement_id", "header_names": ["요구사항ID"], "default_column": 7},
            {"field": "requirement_name", "header_names": ["요구사항명"], "default_column": 8},
            {"field": "description", "header_names": ["기능/비기능요구사항", "기능/비기능요구사항(취합정리 5/25)", "요구사항내용"], "default_column": 9},
            {"field": "user_auth_requirement", "header_names": ["사용자권한요구사항"], "optional": True, "default_column": 10},
            {"field": "request_dept", "header_names": ["의뢰부서"], "optional": True, "default_column": 11},
            {"field": "owner_team", "header_names": ["요구사항처리담당팀"], "optional": True, "default_column": 12},
            {"field": "review_status", "header_names": ["검토 상태", "검토상태"], "optional": True, "default_column": 13},
            {"field": "note", "header_names": ["검토의견"], "optional": True, "default_column": 14},
            {"field": "change_requirement_id", "header_names": ["변경요구사항ID"], "optional": True, "default_column": 15},
            {"field": "change_date", "header_names": ["변경일"], "optional": True, "default_column": 16},
            {"field": "source_doc", "header_names": ["근기문서", "근거문서"], "optional": True, "default_column": 17},
            {"field": "ace", "header_names": ["ACE"], "optional": True, "default_column": 18},
        ]
        merged = list(columns or [])
        seen_fields = {str(column.get("field")) for column in merged if isinstance(column, dict)}
        seen_field_parts = {
            part.strip()
            for field in seen_fields
            for part in field.split("|")
            if part.strip()
        }
        for column in legacy_columns:
            field_parts = {
                part.strip()
                for part in str(column["field"]).split("|")
                if part.strip()
            }
            if column["field"] not in seen_fields and not field_parts.intersection(seen_field_parts):
                merged.append(column)
        return merged

    def _requirement_source(self, item: dict[str, Any]) -> dict[str, Any]:
        metadata = item.get("metadata") or {}
        source_doc = (
            metadata.get("source_doc")
            or item.get("source_doc")
            or metadata.get("source_file_name")
            or item.get("source_file_name")
            or item.get("source_document_id")
            or metadata.get("source_document_id")
            or ""
        )
        return {
            **metadata,
            **item,
            "work": metadata.get("work") or metadata.get("domain") or "",
            "section_category": metadata.get("section_category") or metadata.get("domain") or "",
            "requirement_id": item.get("requirement_id") or metadata.get("requirement_id", ""),
            "requirement_name": metadata.get("requirement_name") or item.get("title", ""),
            "description": metadata.get("description") or item.get("description", ""),
            "source_doc": source_doc,
            "note": metadata.get("note") or "",
        }

    def _unit_test_source(self, item: dict[str, Any]) -> dict[str, Any]:
        metadata = item.get("metadata") or {}
        test_content = item.get("test_content")
        if test_content in (None, ""):
            test_content = metadata.get("test_content", "")
        return {
            **metadata,
            **item,
            "new_category": metadata.get("new_category") or "신규",
            "screen_batch_type": metadata.get("screen_batch_type") or "화면",
            "test_case_id": item.get("test_case_id") or metadata.get("test_case_id", ""),
            "test_case_name": item.get("test_case_name") or metadata.get("test_case_name", ""),
            "requirement_id": item.get("requirement_id") or metadata.get("requirement_id", ""),
            "requirement_name": item.get("requirement_name") or metadata.get("requirement_name", ""),
            "scenario_id": item.get("scenario_id") or metadata.get("scenario_id", ""),
            "test_content": self._unit_test_content_text(test_content),
            "author": metadata.get("author") or "",
        }

    def _unit_test_content_text(self, value: Any) -> str:
        text = (
            str(value or "")
            .replace("\\n", "\n")
            .replace("\r\n", "\n")
            .replace("\r", "\n")
        )
        text = text.strip()
        if not text:
            return " "
        if "\n" not in text:
            return text

        lines: list[str] = []
        for line in text.split("\n"):
            cleaned = line.strip()
            if not cleaned:
                continue
            cleaned = re.sub(r"^\s*(?:[-*•]+|\d+[.)])\s*", "", cleaned).strip()
            if cleaned:
                lines.append(cleaned)
        if len(lines) <= 1:
            return lines[0] if lines else text
        return "\n".join(f"{index}. {line}" for index, line in enumerate(lines, start=1))

    def _format_unit_test_content_column(
        self,
        ws: Any,
        data_mapper: dict[str, Any],
        columns: list[dict[str, Any]],
        item_count: int,
    ) -> None:
        from copy import copy
        from openpyxl.cell.cell import MergedCell

        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))
        headers = self._header_map(ws, header_row=header_row)
        column = None
        for column_mapper in columns:
            if column_mapper.get("field") == "test_content":
                column = self._find_column(headers, column_mapper)
                break
        if column is None:
            return
        for row in range(start_row, start_row + max(item_count, 0)):
            cell = ws.cell(row=row, column=column)
            if isinstance(cell, MergedCell):
                continue
            alignment = copy(cell.alignment)
            alignment.wrap_text = True
            alignment.vertical = alignment.vertical or "top"
            cell.alignment = alignment

    def _resolve_sheet_name(
        self,
        workbook: Any,
        preferred_name: str,
        fallback_names: list[str],
    ) -> str:
        if preferred_name in workbook.sheetnames:
            return preferred_name
        for fallback_name in fallback_names:
            if fallback_name in workbook.sheetnames:
                return fallback_name
        return workbook.sheetnames[-1]

    def _normalize_wbs_hierarchy(self, tasks: list[dict[str, Any]], project_name: str = "프로젝트명") -> list[dict[str, Any]]:
        """Ensure WBS rows always have display IDs like 0, 1, 1.1, 1.1.1.

        The agent normally assigns ``wbs_id`` already, but export can receive
        older artifacts or template rows that only have ``level``.  Recompute the
        display ID immediately before Excel writing so the S3 file always uses
        the expected hierarchy format.
        """
        normalized: list[dict[str, Any]] = []
        counters: list[int] = []
        for raw in tasks or []:
            task = dict(raw or {})
            metadata = dict(task.get("metadata") or {})
            try:
                level = int(str(metadata.get("level", task.get("level", "0"))).strip() or "0")
            except ValueError:
                level = 0
            level = max(level, 0)
            explicit_id = str(metadata.get("wbs_id") or task.get("wbs_id") or task.get("id") or "").strip()
            if explicit_id:
                wbs_id = explicit_id
            elif level == 0:
                counters = [0]
                wbs_id = "0"
                if not str(task.get("name") or "").strip() or str(task.get("name")).strip() in {"프로젝트", "프로젝트명"}:
                    task["name"] = project_name or "프로젝트명"
            else:
                while len(counters) <= level:
                    counters.append(0)
                counters = counters[: level + 1]
                counters[level] += 1
                for idx in range(1, level):
                    if counters[idx] == 0:
                        counters[idx] = 1
                wbs_id = ".".join(str(counters[idx]) for idx in range(1, level + 1))
            metadata["level"] = str(level)
            metadata["wbs_id"] = wbs_id
            metadata["id"] = wbs_id
            task["level"] = str(level)
            task["wbs_id"] = wbs_id
            task["id"] = wbs_id
            normalized.append(task)
        return normalized

    def _wbs_source(self, task: dict[str, Any]) -> dict[str, Any]:
        metadata = task.get("metadata") or {}
        return {
            **metadata,
            **task,
            "id": task.get("wbs_id") or metadata.get("wbs_id") or task.get("task_id", ""),
            "wbs_name": task.get("name", ""),
            "deliverable": metadata.get("deliverable", ""),
            "planned_start_date": (
                task.get("planned_start_date")
                or metadata.get("planned_start_date")
                or metadata.get("start_date")
                or ""
            ),
            "planned_end_date": (
                task.get("planned_end_date")
                or metadata.get("planned_end_date")
                or metadata.get("end_date")
                or ""
            ),
            "start_date": (
                task.get("start_date")
                or metadata.get("start_date")
                or task.get("planned_start_date")
                or metadata.get("planned_start_date")
                or ""
            ),
            "end_date": (
                task.get("end_date")
                or metadata.get("end_date")
                or task.get("planned_end_date")
                or metadata.get("planned_end_date")
                or ""
            ),
        }

    def _apply_wbs_schedule_defaults(
        self,
        tasks: list[dict[str, Any]],
        result_json: dict[str, Any],
    ) -> None:
        """Backfill schedule columns so the Excel export never leaves them blank."""
        artifact_metadata = result_json.get("metadata") or {}
        today_text = date.today().strftime("%Y.%m.%d")
        project_start_date = str(artifact_metadata.get("project_start_date") or today_text).strip()
        project_end_date = str(artifact_metadata.get("project_end_date") or today_text).strip()

        for task in tasks:
            metadata = task.setdefault("metadata", {})
            if not isinstance(metadata, dict):
                metadata = {}
                task["metadata"] = metadata
            try:
                level = int(str(metadata.get("level", task.get("level", "0"))).strip() or "0")
            except ValueError:
                level = 0

            start_date_value = str(
                task.get("start_date")
                or metadata.get("start_date")
                or task.get("planned_start_date")
                or metadata.get("planned_start_date")
                or project_start_date
                or today_text
            ).strip()
            end_date_value = str(
                task.get("end_date")
                or metadata.get("end_date")
                or task.get("planned_end_date")
                or metadata.get("planned_end_date")
                or project_end_date
                or today_text
            ).strip()

            if start_date_value:
                metadata["start_date"] = start_date_value
                task["start_date"] = start_date_value
                metadata.setdefault("planned_start_date", start_date_value)
                task.setdefault("planned_start_date", start_date_value)
            if end_date_value:
                metadata["end_date"] = end_date_value
                task["end_date"] = end_date_value
                metadata.setdefault("planned_end_date", end_date_value)
                task.setdefault("planned_end_date", end_date_value)

            explicit_worker = str(metadata.get("worker") or task.get("worker") or "").strip()
            if level > 1 or explicit_worker:
                metadata["worker"] = explicit_worker or "작업자"
                task["worker"] = metadata["worker"]
            else:
                metadata.pop("worker", None)
                task.pop("worker", None)

    def _write_wbs_required_schedule_fields(
        self,
        ws: Any,
        tasks: list[dict[str, Any]],
        data_mapper: dict[str, Any],
    ) -> None:
        """Force-fill WBS schedule/worker columns using the actual workbook."""
        from openpyxl.cell.cell import MergedCell

        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))
        headers = self._header_map(ws, header_row=header_row)
        column_by_header = {
            str(header).replace("\n", "").strip(): column
            for header, column in headers.items()
        }
        fallback_columns = {
            "시작예정일": 5,
            "종료예정일": 6,
            "작업자": 7,
        }
        forced_fields = [
            ("시작예정일", "start_date|planned_start_date"),
            ("종료예정일", "end_date|planned_end_date"),
            ("작업자", "worker|assignee|owner"),
        ]

        for row_offset, task in enumerate(tasks or []):
            source = self._wbs_source(task)
            excel_row = start_row + row_offset
            for header, field_expr in forced_fields:
                column = column_by_header.get(header) or fallback_columns.get(header)
                if column is None:
                    continue
                cell = ws.cell(row=excel_row, column=column)
                if isinstance(cell, MergedCell):
                    continue
                cell.value = get_value(source, field_expr)

    def _apply_wbs_visual_formatting(
        self,
        ws: Any,
        tasks: list[dict[str, Any]],
        data_mapper: dict[str, Any],
    ) -> None:
        """Apply WBS indentation and section highlight fills."""
        from copy import copy
        from openpyxl.cell.cell import MergedCell
        from openpyxl.styles import Alignment, Color, PatternFill

        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))
        headers = self._header_map(ws, header_row=header_row)
        name_column = None
        for candidate in ("WBS명", "wbs_name"):
            if candidate in headers:
                name_column = headers[candidate]
                break
        if name_column is None:
            name_column = 4

        orange_fill = PatternFill(fill_type="solid")
        orange_fill.fgColor = Color(theme=5, tint=0.7999816888943144)
        yellow_fill = PatternFill(fill_type="solid", fgColor="FFFFFF00")
        orange_names = {"관리영역", "개발영역"}
        yellow_names = {"요구사항정의", "분석", "설계", "구현", "테스트", "이행", "안정화"}

        for row_offset, task in enumerate(tasks or []):
            excel_row = start_row + row_offset
            source = self._wbs_source(task)
            level_raw = str(source.get("level") or task.get("level") or "0").strip()
            try:
                level = max(int(level_raw), 0)
            except ValueError:
                level = 0
            indent_level = level if level <= 1 else level + 1

            name_cell = ws.cell(row=excel_row, column=name_column)
            if not isinstance(name_cell, MergedCell):
                alignment = copy(name_cell.alignment)
                alignment = alignment or Alignment()
                alignment.indent = float(indent_level)
                alignment.wrap_text = True
                name_cell.alignment = alignment

            row_name = str(source.get("wbs_name") or task.get("name") or "").strip()
            if row_name in orange_names:
                fill = orange_fill
            elif row_name in yellow_names:
                fill = yellow_fill
            else:
                fill = None

            if fill is None:
                continue

            for column in range(1, min(12, ws.max_column) + 1):
                cell = ws.cell(row=excel_row, column=column)
                if isinstance(cell, MergedCell):
                    continue
                cell.fill = fill

    def _screen_source(self, screen: dict[str, Any]) -> dict[str, Any]:
        metadata = screen.get("metadata") or {}
        source_requirement_ids = screen.get("source_requirement_ids") or []
        return {
            **metadata,
            **screen,
            "screen_name": screen.get("name", ""),
            "screen_no": screen.get("screen_id", ""),
            "requirement_id": metadata.get("requirement_id") or ", ".join(source_requirement_ids),
            "requirement_name": metadata.get("requirement_name") or screen.get("name", ""),
            "description": metadata.get("description") or screen.get("description", ""),
            "display_items": metadata.get("display_items", []),
        }

    def _common_screen_placeholder_values(self, mapper: dict[str, Any], context: dict[str, str]) -> dict[str, str]:
        values = build_placeholder_values(get_nested(mapper, "placeholder_slides", "common", default={}) or {}, context=context)
        project_name = context.get("project_name") or "프로젝트명"
        values.update({
            "{프로젝트명}": project_name,
            "{ProjectName}": project_name,
            "{PROJECT_NAME}": project_name,
            "{작성일}": context.get("today", ""),
            "{작성자}": context.get("author", ""),
            "{작성자명}": context.get("author", ""),
        })
        return values

    def _screen_placeholder_values(self, source: dict[str, Any], context: dict[str, str]) -> dict[str, str]:
        requirement_id = str(source.get("requirement_id") or "")
        screen_id = str(source.get("screen_id") or source.get("screen_no") or "")
        screen_name = str(source.get("screen_name") or source.get("name") or "")
        requirement_name = str(source.get("requirement_name") or screen_name)
        description = str(source.get("description") or "")
        return {
            "{프로젝트명}": context.get("project_name") or "프로젝트명",
            "{작성일}": context.get("today", ""),
            "{작성자}": context.get("author", ""),
            "{작성자명}": context.get("author", ""),
            "{요구사항ID}": requirement_id,
            "{요구사항명}": requirement_name,
            "{요구사항내용}": description,
            "{Description}": description,
            "{DESCRIPTION}": description,
            "{description}": description,
            "{화면ID}": screen_id,
            "{화면번호}": screen_id,
            "{화면명}": screen_name,
            "{서브시스템명}": str(source.get("biz_requirement_name") or source.get("domain") or ""),
            "{메뉴위치}": str(source.get("domain") or source.get("biz_requirement_name") or ""),
        }

    def _replace_placeholders_in_presentation(self, prs: Any, values: dict[str, str]) -> None:
        for slide in prs.slides:
            self._replace_placeholders_in_slide(slide, values)

    def _apply_placeholder_sheets(self, wb: Any, placeholder_sheets: list[dict[str, Any]], context: dict[str, str]) -> None:
        from openpyxl.cell.cell import MergedCell
        for sheet_mapper in placeholder_sheets:
            sheet_name = sheet_mapper.get("sheet_name")
            if not sheet_name or sheet_name not in wb.sheetnames:
                continue
            values = build_placeholder_values(sheet_mapper.get("placeholders", {}), source=None, context=context)
            ws = wb[sheet_name]
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell, MergedCell) or not isinstance(cell.value, str):
                        continue
                    text = cell.value
                    for key, value in values.items():
                        text = text.replace(key, value or "")
                    cell.value = text

    def _header_map(self, ws: Any, header_row: int = 1) -> dict[str, int]:
        headers: dict[str, int] = {}
        for cell in ws[header_row]:
            if cell.value is None:
                continue
            key = str(cell.value).replace("\n", "").strip()
            if key and key not in headers:
                headers[key] = cell.column
        return headers

    def _find_column(self, headers: dict[str, int], column_mapper: dict[str, Any]) -> int | None:
        if column_mapper.get("column") is not None:
            return int(column_mapper["column"])
        normalized_headers = {name.replace("\n", "").strip(): col for name, col in headers.items()}
        for name in column_mapper.get("header_names") or []:
            key = str(name).replace("\n", "").strip()
            if key in normalized_headers:
                return normalized_headers[key]
        if column_mapper.get("default_column") is not None:
            return int(column_mapper["default_column"])
        if column_mapper.get("optional"):
            return None
        return None

    def _write_mapped_rows(self, ws: Any, items: list[dict[str, Any]], data_mapper: dict[str, Any], columns: list[dict[str, Any]], context: dict[str, str], source_builder: Any) -> None:
        from copy import copy
        from openpyxl.cell.cell import MergedCell
        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))

        # Requirement templates often contain pre-merged sample rows below the
        # header. If we write generated rows into those merged ranges, only the
        # first row receives values and following rows look blank. Unmerge only
        # the data area, keeping cover/history sheet formatting intact.
        for merged_range in list(ws.merged_cells.ranges):
            if merged_range.max_row >= start_row:
                ws.unmerge_cells(str(merged_range))

        headers = self._header_map(ws, header_row=header_row)
        resolved_columns = []
        for column_mapper in columns:
            col_no = self._find_column(headers, column_mapper)
            if col_no is not None:
                resolved_columns.append((col_no, column_mapper))

        # Ensure there are enough physical rows and copy the first data row style
        # to newly created rows so template formatting remains close to sample_0605.
        required_max_row = start_row + max(len(items), 1) - 1
        if ws.max_row < required_max_row and data_mapper.get("copy_template_row_style", True):
            template_row = start_row if ws.max_row >= start_row else header_row
            for row_idx in range(ws.max_row + 1, required_max_row + 1):
                for col_idx in range(1, ws.max_column + 1):
                    src = ws.cell(row=template_row, column=col_idx)
                    dst = ws.cell(row=row_idx, column=col_idx)
                    if src.has_style:
                        dst._style = copy(src._style)
                    dst.number_format = src.number_format
                    dst.alignment = copy(src.alignment)
                    dst.font = copy(src.font)
                    dst.fill = copy(src.fill)
                    dst.border = copy(src.border)

        if data_mapper.get("clear_existing_rows", True):
            for row in range(start_row, ws.max_row + 1):
                for col_no, _ in resolved_columns:
                    cell = ws.cell(row=row, column=col_no)
                    if not isinstance(cell, MergedCell):
                        cell.value = None

        for row_offset, item in enumerate(items):
            excel_row = start_row + row_offset
            row_number = row_offset + 1
            source = source_builder(item)
            for col_no, column_mapper in resolved_columns:
                cell = ws.cell(row=excel_row, column=col_no)
                if isinstance(cell, MergedCell):
                    continue
                cell.value = get_value(source, column_mapper.get("field", ""), context=context, row_number=row_number)

    def _write_requirement_reference_fields(
        self,
        ws: Any,
        items: list[dict[str, Any]],
        data_mapper: dict[str, Any],
    ) -> None:
        """Force reference-specific requirement columns after mapper writes.

        Runtime can prefer an S3-hosted output_mapper.json. This final pass uses
        the actual template headers so stale mapper aliases cannot overwrite
        업무/구분 with Biz요건명 or copy description into 검토의견.
        """
        from openpyxl.cell.cell import MergedCell

        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))
        headers = self._header_map(ws, header_row=header_row)
        column_by_header = {
            str(header).replace("\n", "").strip(): column
            for header, column in headers.items()
        }
        default_requirement_columns = (
            {
                "업무": 1,
                "구분": 2,
                "요구사항구분": 5,
                "기능/비기능요구사항": 11,
                "검토의견": 16,
            }
            if ws.title == "요구사항명세서"
            else {}
        )
        forced_fields = [
            ("업무", "work|domain"),
            ("구분", "section_category|domain"),
            ("요구사항구분", "category"),
            ("기능/비기능요구사항", "description"),
            ("검토의견", "note"),
        ]
        for row_offset, item in enumerate(items or []):
            source = self._requirement_source(item)
            excel_row = start_row + row_offset
            for header, field_expr in forced_fields:
                column = default_requirement_columns.get(header) or column_by_header.get(header)
                if column is None:
                    continue
                cell = ws.cell(row=excel_row, column=column)
                if isinstance(cell, MergedCell):
                    continue
                cell.value = get_value(source, field_expr) if header != "검토의견" else ""

    def _write_wbs_display_ids(
        self,
        ws: Any,
        tasks: list[dict[str, Any]],
        data_mapper: dict[str, Any],
        columns: list[dict[str, Any]],
    ) -> None:
        """Force-fill the WBS display ID column after template mapping."""
        from openpyxl.cell.cell import MergedCell

        header_row = int(data_mapper.get("header_row", 1))
        start_row = int(data_mapper.get("start_row", 2))
        headers = self._header_map(ws, header_row=header_row)
        id_column = None
        for column_mapper in columns:
            if column_mapper.get("field") in {"id", "wbs_id"}:
                id_column = self._find_column(headers, column_mapper)
                break
        if id_column is None:
            id_column = 3

        for row_offset, task in enumerate(tasks):
            source = self._wbs_source(task)
            display_id = source.get("id") or source.get("wbs_id")
            if display_id in (None, ""):
                continue
            cell = ws.cell(row=start_row + row_offset, column=id_column)
            if not isinstance(cell, MergedCell):
                cell.value = display_id

    def _replace_common_slide_placeholders(self, prs: Any, mapper: dict[str, Any], context: dict[str, str]) -> None:
        common_values = build_placeholder_values(get_nested(mapper, "placeholder_slides", "common", default={}) or {}, context=context)
        for slide_index in mapper.get("common_slide_indices", [0, 1]):
            if slide_index < len(prs.slides):
                self._replace_placeholders_in_slide(prs.slides[slide_index], common_values)

    def _replace_placeholders_in_slide(self, slide: Any, values: dict[str, str]) -> None:
        def replace_text(value: str) -> str:
            text = str(value or "")
            for key, replacement in values.items():
                text = text.replace(key, replacement or "")
            return text

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                self._replace_text_frame_preserving_style(
                    shape.text_frame,
                    replace_text,
                )
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        self._replace_text_frame_preserving_style(
                            cell.text_frame,
                            replace_text,
                        )

    def _replace_text_frame_preserving_style(self, text_frame: Any, replace_text: Any) -> None:
        for paragraph in text_frame.paragraphs:
            runs = list(paragraph.runs)
            if runs:
                original = "".join(run.text or "" for run in runs)
                replaced = replace_text(original)
                if replaced != original:
                    runs[0].text = replaced
                    for run in runs[1:]:
                        run.text = ""
                    continue
                for run in runs:
                    run.text = replace_text(run.text)
                continue

            original = paragraph.text or ""
            replaced = replace_text(original)
            if replaced != original:
                paragraph.text = replaced

    def _set_text_frame_text_preserving_style(
        self,
        text_frame: Any,
        text: str,
        style_source_text_frame: Any = None,
    ) -> None:
        paragraphs = list(text_frame.paragraphs)
        if not paragraphs:
            text_frame.text = text
            return
        paragraph = paragraphs[0]
        runs = list(paragraph.runs)
        source_run = self._first_run(style_source_text_frame)
        if runs:
            if source_run is not None:
                self._copy_run_font(source_run, runs[0])
            runs[0].text = text
            for run in runs[1:]:
                run.text = ""
        else:
            run = paragraph.add_run()
            if source_run is not None:
                self._copy_run_font(source_run, run)
            run.text = text
        for extra_paragraph in paragraphs[1:]:
            for run in extra_paragraph.runs:
                run.text = ""

    def _first_paragraph_with_run(self, text_frame: Any) -> Any:
        if text_frame is None:
            return None
        for paragraph in text_frame.paragraphs:
            if list(paragraph.runs):
                return paragraph
        return None

    def _first_run(self, text_frame: Any) -> Any:
        if text_frame is None:
            return None
        for paragraph in text_frame.paragraphs:
            runs = list(paragraph.runs)
            if runs:
                return runs[0]
        return None

    def _copy_run_font(self, source_run: Any, target_run: Any) -> None:
        source_font = source_run.font
        target_font = target_run.font
        target_font.name = source_font.name
        target_font.size = source_font.size
        target_font.bold = source_font.bold
        target_font.italic = source_font.italic
        target_font.underline = source_font.underline
        if source_font.color is not None:
            try:
                if source_font.color.rgb is not None:
                    target_font.color.rgb = source_font.color.rgb
            except AttributeError:
                pass

    def _fill_description_table(self, slide: Any, source: dict[str, Any], table_mapper: dict[str, Any]) -> bool:
        start_row = int(table_mapper.get("start_row", 1))
        target_column = int(table_mapper.get("target_column", 1))
        clear_rows = bool(table_mapper.get("clear_rows_before_fill", True))
        header_text = str(table_mapper.get("header_text") or "").strip()
        min_rows = int(table_mapper.get("min_rows", 0))
        min_columns = int(table_mapper.get("min_columns", 0))
        description_lines = self._screen_description_lines(source)
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if min_rows and len(table.rows) < min_rows:
                continue
            if min_columns and len(table.columns) < min_columns:
                continue
            if len(table.columns) <= target_column:
                continue
            if header_text and not self._table_contains_text(table, header_text):
                continue
            default_style_cell = self._description_table_style_cell(
                table,
                start_row=start_row,
                target_column=target_column,
            )
            if clear_rows:
                for row_idx in range(start_row, len(table.rows)):
                    try:
                        style_cell = self._row_style_cell(
                            table,
                            row_idx=row_idx,
                            target_column=target_column,
                            default_style_cell=default_style_cell,
                        )
                        self._set_text_frame_text_preserving_style(
                            table.cell(row_idx, target_column).text_frame,
                            "",
                            style_cell.text_frame if style_cell is not None else None,
                        )
                    except Exception:
                        pass
            for offset, description in enumerate(description_lines):
                row_idx = start_row + offset
                if row_idx >= len(table.rows):
                    break
                style_cell = self._row_style_cell(
                    table,
                    row_idx=row_idx,
                    target_column=target_column,
                    default_style_cell=default_style_cell,
                )
                self._set_text_frame_text_preserving_style(
                    table.cell(row_idx, target_column).text_frame,
                    description,
                    style_cell.text_frame if style_cell is not None else None,
                )
            return True
        return False

    def _screen_description_lines(self, source: dict[str, Any]) -> list[str]:
        description = str(source.get("description") or "").strip()
        if not description:
            return []
        lines = [
            line.strip()
            for line in description.replace("\r\n", "\n").replace("\r", "\n").split("\n")
            if line.strip()
        ]
        if not lines:
            return []
        return lines

    def _description_table_style_cell(
        self,
        table: Any,
        *,
        start_row: int,
        target_column: int,
    ) -> Any:
        for row_idx in range(start_row, len(table.rows)):
            cell = table.cell(row_idx, target_column)
            if self._first_run(cell.text_frame) is not None:
                return cell
        for row_idx in range(start_row, len(table.rows)):
            for col_idx in range(len(table.columns)):
                if col_idx == target_column:
                    continue
                cell = table.cell(row_idx, col_idx)
                if self._first_run(cell.text_frame) is not None:
                    return cell
        for row in table.rows:
            for cell in row.cells:
                if self._first_run(cell.text_frame) is not None:
                    return cell
        return None

    def _row_style_cell(
        self,
        table: Any,
        *,
        row_idx: int,
        target_column: int,
        default_style_cell: Any,
    ) -> Any:
        target_cell = table.cell(row_idx, target_column)
        if self._first_run(target_cell.text_frame) is not None:
            return target_cell
        for col_idx in range(len(table.columns)):
            if col_idx == target_column:
                continue
            cell = table.cell(row_idx, col_idx)
            if self._first_run(cell.text_frame) is not None:
                return cell
        return default_style_cell

    def _table_contains_text(self, table: Any, expected_text: str) -> bool:
        expected = expected_text.replace("\n", "").strip()
        if not expected:
            return True
        for row in table.rows:
            for cell in row.cells:
                text = str(cell.text or "").replace("\n", "").strip()
                if text == expected:
                    return True
        return False

    def _add_description_fallback(self, slide: Any, source: dict[str, Any]) -> None:
        from pptx.util import Inches, Pt
        text = "\n".join([
            f"요구사항ID: {source.get('requirement_id', '')}",
            f"요구사항명: {source.get('requirement_name') or source.get('screen_name', '')}",
            f"Description: {source.get('description', '')}",
        ]).strip()
        box = slide.shapes.add_textbox(Inches(8.0), Inches(1.15), Inches(4.7), Inches(4.8))
        box.text_frame.text = text
        for paragraph in box.text_frame.paragraphs:
            paragraph.font.size = Pt(11)

    def _get_blank_slide_layout(self, prs: Any) -> Any:
        """Return a usable blank-like slide layout without assuming layout index 6 exists."""
        layouts = list(prs.slide_layouts)
        if not layouts:
            raise ValueError("PPTX template has no slide layouts")
        blank_candidates = [
            layout for layout in layouts
            if "blank" in (getattr(layout, "name", "") or "").lower()
            or "빈" in (getattr(layout, "name", "") or "")
        ]
        if blank_candidates:
            return blank_candidates[0]
        return min(layouts, key=lambda layout: len(getattr(layout, "placeholders", [])))

    def _duplicate_slide(self, prs: Any, template_slide_index: int) -> Any:
        from copy import deepcopy

        source = prs.slides[template_slide_index]
        dest = prs.slides.add_slide(source.slide_layout)
        self._clear_slide_shapes(dest)
        self._copy_slide_background(source, dest)

        for shape in source.shapes:
            dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

        self._copy_slide_relationships(source, dest)
        return dest

    def _copy_slide_background(self, source_slide: Any, dest_slide: Any) -> None:
        from copy import deepcopy

        source_bg = getattr(getattr(source_slide, "_element", None), "cSld", None)
        source_bg = getattr(source_bg, "bg", None)
        if source_bg is None:
            return
        dest_c_sld = getattr(getattr(dest_slide, "_element", None), "cSld", None)
        if dest_c_sld is None:
            return
        dest_bg = getattr(dest_c_sld, "bg", None)
        if dest_bg is not None:
            dest_c_sld.remove(dest_bg)
        dest_c_sld.insert(0, deepcopy(source_bg))

    def _copy_slide_relationships(self, source_slide: Any, dest_slide: Any) -> None:
        skip_reltypes = {
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        }
        for rel in getattr(source_slide.part.rels, "_rels", {}).values():
            if rel.reltype in skip_reltypes:
                continue
            target = rel.target_ref if rel.is_external else rel.target_part
            dest_slide.part.relate_to(target, rel.reltype, is_external=rel.is_external)

    def _clear_slide_shapes(self, slide: Any) -> None:
        for shape in list(slide.shapes):
            shape.element.getparent().remove(shape.element)

    def _delete_slide(self, prs: Any, slide_index: int) -> None:
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        if 0 <= slide_index < len(slides):
            xml_slides.remove(slides[slide_index])


    def _style_sheet(self, ws: Any, widths: list[int]) -> None:
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter

        header_fill = PatternFill(fill_type="solid", fgColor="D9EAF7")
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        for idx, width in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(idx)].width = width
        ws.freeze_panes = "A2"
        ws.auto_filter.ref = ws.dimensions

    def _save_workbook(self, workbook: Any) -> bytes:
        bio = BytesIO()
        workbook.save(bio)
        return bio.getvalue()


artifact_export_service = ArtifactExportService()
