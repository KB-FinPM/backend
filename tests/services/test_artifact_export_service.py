# EN: Tests for generated artifact file export behavior.
# KO: ?앹꽦 ?곗텧臾??뚯씪 export ?숈옉 ?뚯뒪?몄엯?덈떎.

from io import BytesIO

import pytest
from openpyxl import load_workbook
from pptx import Presentation

import app.db.base  # noqa: F401
from app.schemas.artifact import ArtifactType
from app.services.artifact_export_service import ArtifactExportService


def test_wbs_export_fills_hierarchical_display_ids(monkeypatch) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    levels = [0, 1, 2, 3, 4, 4, 3, 4, 4, 3, 4, 4, 3, 4, 4]
    expected_ids = [
        "0",
        "1",
        "1.1",
        "1.1.1",
        "1.1.1.1",
        "1.1.1.2",
        "1.1.2",
        "1.1.2.1",
        "1.1.2.2",
        "1.1.3",
        "1.1.3.1",
        "1.1.3.2",
        "1.1.4",
        "1.1.4.1",
        "1.1.4.2",
    ]
    result_json = {
        "artifact_type": "WBS",
        "tasks": [
            {
                "task_id": f"WBS-{index:03d}",
                "name": f"Task {index}",
                "metadata": {"level": str(level)},
            }
            for index, level in enumerate(levels, start=1)
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_wbs_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=False)
    sheet = workbook["WBS"]

    actual_ids = [sheet.cell(row=row, column=3).value for row in range(2, 17)]

    assert actual_ids == expected_ids


def test_wbs_export_orders_rows_by_wbs_hierarchy(monkeypatch) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "WBS",
        "tasks": [
            {
                "task_id": "WBS-004",
                "name": "이행 상세",
                "metadata": {"level": "4", "wbs_id": "3.6.2.3"},
            },
            {
                "task_id": "WBS-002",
                "name": "요구사항 정의 상세 2",
                "metadata": {"level": "4", "wbs_id": "3.1.2"},
            },
            {
                "task_id": "WBS-003",
                "name": "설계 상세",
                "metadata": {"level": "4", "wbs_id": "3.3.5"},
            },
            {
                "task_id": "WBS-001",
                "name": "요구사항 정의 상세 1",
                "metadata": {"level": "4", "wbs_id": "3.1.1"},
            },
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_wbs_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=False)
    sheet = workbook["WBS"]

    actual_ids = [sheet.cell(row=row, column=3).value for row in range(2, 6)]

    assert actual_ids == ["3.1.1", "3.1.2", "3.3.5", "3.6.2.3"]


def test_wbs_export_fills_planned_date_columns(monkeypatch) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "WBS",
        "tasks": [
            {
                "task_id": "WBS-001",
                "name": "Project",
                "planned_start_date": "2024-01-10",
                "planned_end_date": "2024-07-09",
                "metadata": {
                    "level": "0",
                    "assignee": "Planner",
                    "status": "Ready",
                },
            }
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_wbs_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    sheet = workbook["WBS"]

    assert sheet.cell(row=2, column=5).value == "2024-01-10"
    assert sheet.cell(row=2, column=6).value == "2024-07-09"
    assert sheet.cell(row=2, column=7).value == "Planner"
    assert sheet.cell(row=2, column=12).value == "Ready"


def test_wbs_export_uses_metadata_date_aliases(monkeypatch) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "WBS",
        "tasks": [
            {
                "task_id": "WBS-001",
                "name": "Project",
                "metadata": {
                    "level": "0",
                    "start_date": "2024-02-01",
                    "end_date": "2024-02-29",
                    "worker": "Worker A",
                    "work_status": "In Progress",
                },
            }
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_wbs_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    sheet = workbook["WBS"]

    assert sheet.cell(row=2, column=5).value == "2024-02-01"
    assert sheet.cell(row=2, column=6).value == "2024-02-29"
    assert sheet.cell(row=2, column=7).value == "Worker A"
    assert sheet.cell(row=2, column=12).value == "In Progress"


def test_wbs_export_allows_missing_optional_schedule_values(monkeypatch) -> None:
    class FixedDate:
        @classmethod
        def today(cls):
            from datetime import date

            return date(2026, 6, 10)

    monkeypatch.setattr("app.services.artifact_export_service.date", FixedDate)
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "WBS",
        "tasks": [
            {
                "task_id": "WBS-001",
                "name": "Project",
                "metadata": {
                    "level": "0",
                    "deliverable": "Plan",
                },
            }
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_wbs_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    sheet = workbook["WBS"]

    assert sheet.cell(row=2, column=2).value == "0"
    assert sheet.cell(row=2, column=3).value == "0"
    assert sheet.cell(row=2, column=4).value == "Project"
    assert sheet.cell(row=2, column=5).value == "2026.06.10"
    assert sheet.cell(row=2, column=6).value == "2026.06.10"
    assert sheet.cell(row=2, column=7).value is None
    assert sheet.cell(row=2, column=8).value == "Plan"
    assert sheet.cell(row=2, column=12).value is None


def test_requirement_export_keeps_work_category_and_empty_review_note(
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "REQUIREMENT_SPEC",
        "requirements": [
            {
                "requirement_id": "BSR-00001",
                "title": "?쒖뒪???꾪궎?띿퀜 ?ㅺ퀎",
                "description": "- ?쒕퉬?ㅻ퀎 ?낅┰?? ?뺤옣 ?좎뿰?깆쓣 媛吏????덈뒗 OCP ?뚮옯??援ъ텞",
                "metadata": {
                    "work": "鍮꾨?硫??꾪궎?띿퀜 ?ъ꽕怨?諛??명봽??援ъ텞",
                    "section_category": "鍮꾨?硫??꾪궎?띿퀜 ?ъ꽕怨?諛??명봽??援ъ텞",
                    "biz_requirement_id": "Biz-0001",
                    "biz_requirement_name": "?꾪궎?띿퀜 ?ㅺ퀎",
                    "category": "湲곕뒫",
                    "creation_stage": "?붽뎄?ы빆?뺤쓽",
                    "status": "?좉퇋",
                    "source": "Construction Requirement Definition",
                    "requirement_name": "?쒖뒪???꾪궎?띿퀜 ?ㅺ퀎",
                    "description": "- ?쒕퉬?ㅻ퀎 ?낅┰?? ?뺤옣 ?좎뿰?깆쓣 媛吏????덈뒗 OCP ?뚮옯??援ъ텞",
                    "note": "",
                },
            }
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_requirement_xlsx(result_json)
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    sheet = workbook["요구사항명세서"]

    assert sheet.cell(row=2, column=1).value == "鍮꾨?硫??꾪궎?띿퀜 ?ъ꽕怨?諛??명봽??援ъ텞"
    assert sheet.cell(row=2, column=2).value == "鍮꾨?硫??꾪궎?띿퀜 ?ъ꽕怨?諛??명봽??援ъ텞"
    assert sheet.cell(row=2, column=5).value == "湲곕뒫"
    assert sheet.cell(row=2, column=11).value == "- ?쒕퉬?ㅻ퀎 ?낅┰?? ?뺤옣 ?좎뿰?깆쓣 媛吏????덈뒗 OCP ?뚮옯??援ъ텞"
    assert sheet.cell(row=2, column=16).value is None


@pytest.mark.anyio
async def test_requirement_export_uses_unprefixed_file_name() -> None:
    class FakeStorageService:
        def __init__(self) -> None:
            self.uploaded_keys: list[str] = []

        async def upload(self, *, file_bytes: bytes, key: str, content_type: str) -> str:
            self.uploaded_keys.append(key)
            return f"mock://{key}"

    storage_service = FakeStorageService()
    result = await ArtifactExportService().export_artifact(
        project_id="PRJ-TEST-001",
        artifact_id="ART-REQ-001",
        artifact_type=ArtifactType.REQUIREMENT_SPEC,
        result_json={
            "artifact_type": "REQUIREMENT_SPEC",
            "requirements": [],
            "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
        },
        project_name="KB Star Banking Process",
        storage_service=storage_service,
    )

    assert result is not None
    assert result.file_name.startswith("[KB Star Banking Process] 요구사항명세서_")
    assert result.file_name.endswith("_v1.xlsx")
    assert storage_service.uploaded_keys[0].endswith(
        f"/REQUIREMENT_SPEC/ART-REQ-001/{result.file_name}"
    )


def test_screen_design_export_creates_pages_with_requirement_descriptions(
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        "util.agent_template_utils.settings.S3_STORAGE_BACKEND",
        "mock",
    )
    result_json = {
        "artifact_type": "SCREEN_DESIGN",
        "screens": [
            {
                "screen_id": "SCR-001",
                "name": "?뚯썝 議고쉶 ?붾㈃",
                "description": "?ъ슜?먮뒗 ?뚯썝 紐⑸줉??議고쉶?????덉뼱???쒕떎.",
                "source_requirement_ids": ["REQ-0001"],
                "metadata": {
                    "requirement_id": "REQ-0001",
                    "requirement_name": "?뚯썝 議고쉶",
                    "description": "?ъ슜?먮뒗 ?뚯썝 紐⑸줉??議고쉶?????덉뼱???쒕떎.",
                    "display_items": [
                        {
                            "item_name": "Description",
                            "description": "?ъ슜?먮뒗 ?뚯썝 紐⑸줉??議고쉶?????덉뼱???쒕떎.",
                        }
                    ],
                },
            },
            {
                "screen_id": "SCR-002",
                "name": "沅뚰븳 愿由??붾㈃",
                "description": "愿由ъ옄???ъ슜??沅뚰븳??蹂寃쏀븷 ???덉뼱???쒕떎.",
                "source_requirement_ids": ["REQ-0002"],
                "metadata": {
                    "requirement_id": "REQ-0002",
                    "requirement_name": "Permission Management",
                    "description": "愿由ъ옄???ъ슜??沅뚰븳??蹂寃쏀븷 ???덉뼱???쒕떎.",
                    "display_items": [
                        {
                            "item_name": "Description",
                            "description": "愿由ъ옄???ъ슜??沅뚰븳??蹂寃쏀븷 ???덉뼱???쒕떎.",
                        }
                    ],
                },
            },
        ],
        "metadata": {"project_id": "PRJ-TEST-001", "author": "Tester"},
    }

    file_bytes = ArtifactExportService()._build_screen_design_pptx(result_json)
    presentation = Presentation(BytesIO(file_bytes))
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
    all_text = "\n".join(texts)

    assert "Tester" in all_text
    assert "REQ-0001" in all_text
    assert "REQ-0002" in all_text
    assert "?ъ슜?먮뒗 ?뚯썝 紐⑸줉??議고쉶?????덉뼱???쒕떎." in all_text
    assert "?묒뾽 ?댁슜???뺤쓽?쒕떎" not in all_text
    assert "寃??議곌굔" not in all_text
    assert "泥섎━ 踰꾪듉" not in all_text
    description_cell, style_cell = _first_description_value_and_style_cells(
        presentation,
    )
    assert description_cell.text == "?ъ슜?먮뒗 ?뚯썝 紐⑸줉??議고쉶?????덉뼱???쒕떎."
    assert "REQ-0001" not in description_cell.text
    assert "?뚯썝 議고쉶" not in description_cell.text
    assert "Description:" not in description_cell.text
    description_run = _first_run(description_cell)
    style_run = _first_run(style_cell)
    assert description_run is not None
    assert style_run is not None
    assert description_run.font.size == style_run.font.size
    assert description_run.font.bold == style_run.font.bold


def _first_description_value_and_style_cells(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if len(table.rows) < 2 or len(table.columns) < 2:
                continue
            if table.cell(0, 0).text.replace("\n", "").strip() != "Description":
                continue
            return table.cell(1, 1), table.cell(1, 0)
    raise AssertionError("Description table not found")


def _first_run(cell):
    for paragraph in cell.text_frame.paragraphs:
        runs = list(paragraph.runs)
        if runs:
            return runs[0]
    return None
